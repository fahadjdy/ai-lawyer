# -*- coding: utf-8 -*-
"""
LexCase — Features Overview deck generator.
Non-technical audience · Simple English · hand-drawn mockups (no screenshots).
Run:  python presentation/build_deck.py
Output: presentation/LexCase-Features-Overview.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ----------------------------------------------------------------------------- palette
def C(h): return RGBColor.from_string(h)

INK        = C('0F172A')   # slate-900
SLATE      = C('475569')   # slate-600
MUTED      = C('94A3B8')   # slate-400
FAINT      = C('CBD5E1')   # slate-300
LINE       = C('E2E8F0')   # slate-200
BG         = C('F8FAFC')   # slate-50
CARD       = C('FFFFFF')
WHITE      = C('FFFFFF')

INDIGO     = C('4F46E5')
INDIGO_DK  = C('3730A3')
INDIGO_LT  = C('EEF2FF')
EMERALD    = C('10B981')
EMERALD_DK = C('047857')
EMERALD_LT = C('ECFDF5')
AMBER      = C('F59E0B')
AMBER_DK   = C('B45309')
AMBER_LT   = C('FFFBEB')
ROSE       = C('F43F5E')
ROSE_DK    = C('BE123C')
ROSE_LT    = C('FFF1F2')
SKY        = C('0EA5E9')
SKY_DK     = C('0369A1')
SKY_LT     = C('F0F9FF')
VIOLET     = C('8B5CF6')
VIOLET_LT  = C('F5F3FF')

FONT  = 'Segoe UI'
FONTB = 'Segoe UI Semibold'
EMO   = 'Segoe UI Emoji'

# ----------------------------------------------------------------------------- setup
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

_page = {'n': 0}


def slide():
    return prs.slides.add_slide(BLANK)


def _noshadow(sp):
    try:
        sp.shadow.inherit = False
    except Exception:
        pass


def rrect(s, x, y, w, h, fill, line=None, lw=1.0, radius=0.12):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    try:
        sp.adjustments[0] = min(0.5, radius / min(w, h))
    except Exception:
        pass
    _noshadow(sp)
    return sp


def rect(s, x, y, w, h, fill, line=None, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    _noshadow(sp)
    return sp


def oval(s, x, y, w, h, fill, line=None, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    _noshadow(sp)
    return sp


def donut(s, x, y, d, fill, thickness=0.22):
    sp = s.shapes.add_shape(MSO_SHAPE.DONUT,
                            Inches(x), Inches(y), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    try:
        sp.adjustments[0] = thickness
    except Exception:
        pass
    _noshadow(sp)
    return sp


def hline(s, x, y, w, color=LINE, weight=1.0):
    ln = s.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    _noshadow(ln)
    return ln


def vline(s, x, y, h, color=LINE, weight=1.0):
    ln = s.shapes.add_connector(2, Inches(x), Inches(y), Inches(x), Inches(y + h))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    _noshadow(ln)
    return ln


def R(t, size=14, color=INK, bold=False, font=FONT, italic=False):
    return {'t': t, 'size': size, 'color': color, 'bold': bold,
            'font': font, 'italic': italic}


def P(runs, align=PP_ALIGN.LEFT, space_after=0, space_before=0, line=None):
    return {'runs': runs, 'align': align, 'space_after': space_after,
            'space_before': space_before, 'line': line}


def tb(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        para.space_after = Pt(p.get('space_after', 0))
        para.space_before = Pt(p.get('space_before', 0))
        if p.get('line'):
            para.line_spacing = p['line']
        for r in p['runs']:
            run = para.add_run(); run.text = r['t']
            f = run.font
            f.size = Pt(r['size']); f.bold = r['bold']; f.italic = r['italic']
            f.color.rgb = r['color']; f.name = r['font']
    return box


def line(s, x, y, w, h, text, size=14, color=INK, bold=False, font=FONT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    return tb(s, x, y, w, h, [P([R(text, size, color, bold, font)], align)], anchor)


def chip(s, x, y, text, fill, fg, size=9.5, padx=0.11, h=0.26, bold=True, font=FONT):
    w = padx * 2 + max(0.18, len(text) * size * 0.0095)
    rrect(s, x, y, w, h, fill, radius=h / 2)
    tb(s, x, y - 0.01, w, h, [P([R(text, size, fg, bold, font)], PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)
    return w


def icon_badge(s, x, y, d, glyph, fill, fg=WHITE, size=15):
    rrect(s, x, y, d, d, fill, radius=d * 0.30)
    tb(s, x, y - 0.01, d, d, [P([R(glyph, size, fg, True, EMO)], PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)


def avatar(s, x, y, d, initials, fill):
    oval(s, x, y, d, d, fill)
    tb(s, x, y - 0.01, d, d,
       [P([R(initials, d * 13, WHITE, True)], PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)


def btn(s, x, y, text, fill=INDIGO, fg=WHITE, size=10, h=0.32, padx=0.16):
    w = padx * 2 + len(text) * size * 0.0105
    rrect(s, x, y, w, h, fill, radius=h / 2)
    tb(s, x, y - 0.01, w, h, [P([R(text, size, fg, True)], PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)
    return w


def footer(s, dark=False):
    _page['n'] += 1
    col = C('64748B') if not dark else C('94A3B8')
    line(s, 0.7, SH - 0.42, 6, 0.3, "LexCase  ·  AI-Assisted Legal Case Management",
         9, col)
    line(s, SW - 1.4, SH - 0.42, 0.7, 0.3, str(_page['n']).zfill(2), 9, col,
         align=PP_ALIGN.RIGHT)


def bg(s, color=BG):
    rect(s, -0.06, -0.06, SW + 0.12, SH + 0.12, color)


# ----------------------------------------------------------------------------- app window frame
def window(s, x, y, w, h, title="lexcase.app"):
    rrect(s, x + 0.05, y + 0.07, w, h, C('E2E8F0'), radius=0.14)   # soft shadow
    rrect(s, x, y, w, h, CARD, line=LINE, lw=1.0, radius=0.14)
    barh = 0.40
    rrect(s, x, y, w, barh + 0.12, C('F1F5F9'), radius=0.14)
    rect(s, x, y + barh - 0.06, w, 0.18, C('F1F5F9'))
    oval(s, x + 0.18, y + 0.14, 0.12, 0.12, ROSE)
    oval(s, x + 0.36, y + 0.14, 0.12, 0.12, AMBER)
    oval(s, x + 0.54, y + 0.14, 0.12, 0.12, EMERALD)
    rrect(s, x + 0.85, y + 0.10, min(3.2, w - 1.6), 0.22, WHITE, line=LINE, radius=0.11)
    line(s, x + 0.98, y + 0.085, 3.0, 0.22, "🔒  " + title, 8.5, MUTED, font=EMO,
         anchor=MSO_ANCHOR.MIDDLE)
    hline(s, x, y + barh + 0.06, w, LINE, 1.0)
    return (x + 0.22, y + barh + 0.22, w - 0.44, h - barh - 0.44)


# ----------------------------------------------------------------------------- feature slide
def feature(num, kicker, kglyph, kcolor, kbg, title, subtitle, bullets, mockup,
            win_title="lexcase.app"):
    s = slide()
    bg(s)
    rect(s, 0, 0, 0.16, SH, kcolor)                       # left accent spine
    # kicker pill
    icon_badge(s, 0.7, 0.62, 0.46, kglyph, kcolor, size=17)
    line(s, 1.28, 0.60, 5, 0.26, kicker.upper(), 11, kcolor, bold=True)
    line(s, 1.28, 0.84, 5, 0.26, "FEATURE " + str(num).zfill(2), 9, MUTED, bold=True)
    # title + subtitle
    tb(s, 0.7, 1.35, 5.05, 1.2,
       [P([R(title, 27, INK, True)], line=1.04)])
    tb(s, 0.7, 2.45, 5.0, 1.1,
       [P([R(subtitle, 13, SLATE)], line=1.18)])
    # bullets
    by = 3.55
    for lead, desc in bullets:
        rrect(s, 0.72, by + 0.03, 0.2, 0.2, kbg, radius=0.06)
        tb(s, 0.74, by + 0.0, 0.18, 0.22,
           [P([R("✓", 10, kcolor, True)], PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        tb(s, 1.08, by - 0.02, 4.7, 0.8,
           [P([R(lead + "  ", 12.5, INK, True), R(desc, 12.5, SLATE)], line=1.12)])
        by += 0.72
    # mockup window on the right
    wx, wy, ww, wh = 6.25, 1.15, 6.35, 5.05
    inner = window(s, wx, wy, ww, wh, win_title)
    mockup(s, *inner)
    footer(s)
    return s


# ============================================================================= MOCKUPS
def mk_dashboard(s, x, y, w, h):
    # KPI row
    kw = (w - 0.36) / 4
    data = [("48", "Active cases", INDIGO, "📂"), ("126", "Clients", SKY, "👥"),
            ("31", "Open tasks", AMBER, "✅"), ("9", "This week", EMERALD, "📅")]
    for i, (num, lab, col, g) in enumerate(data):
        cx = x + i * (kw + 0.12)
        rrect(s, cx, y, kw, 1.05, WHITE, line=LINE, radius=0.10)
        icon_badge(s, cx + 0.12, y + 0.14, 0.34, g, col, size=12)
        line(s, cx + 0.12, y + 0.52, kw, 0.4, num, 20, INK, bold=True)
        line(s, cx + 0.13, y + 0.84, kw, 0.3, lab, 8.5, MUTED)
    # chart card (bars)
    cy = y + 1.25
    ch = h - 1.35
    cardw = w * 0.60
    rrect(s, x, cy, cardw, ch, WHITE, line=LINE, radius=0.10)
    line(s, x + 0.18, cy + 0.14, cardw, 0.3, "Cases & hearings — last 6 months",
         9.5, INK, bold=True)
    base = cy + ch - 0.45
    bw = 0.34
    gap = (cardw - 0.5 - 6 * bw) / 6
    heights = [0.9, 1.3, 1.05, 1.7, 1.45, 2.0]
    for i, hh in enumerate(heights):
        bx = x + 0.32 + i * (bw + gap)
        rrect(s, bx, base - hh, bw * 0.46, hh, C('C7D2FE'), radius=0.04)
        rrect(s, bx + bw * 0.5, base - hh * 0.7, bw * 0.46, hh * 0.7, INDIGO, radius=0.04)
    hline(s, x + 0.2, base, cardw - 0.4, LINE, 1.0)
    # donut card (win rate)
    dx = x + cardw + 0.18
    dw = w - cardw - 0.18
    rrect(s, dx, cy, dw, ch, WHITE, line=LINE, radius=0.10)
    line(s, dx + 0.16, cy + 0.14, dw, 0.3, "Win rate", 9.5, INK, bold=True)
    dd = min(dw - 0.5, ch - 1.1)
    donut(s, dx + (dw - dd) / 2, cy + 0.5, dd, EMERALD, thickness=0.26)
    tb(s, dx, cy + 0.5 + dd / 2 - 0.28, dw, 0.6,
       [P([R("68%", 19, INK, True)], PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    line(s, dx, cy + ch - 0.42, dw, 0.3, "cases won", 8.5, MUTED, align=PP_ALIGN.CENTER)


def mk_ai(s, x, y, w, h):
    icon_badge(s, x, y, 0.40, "✨", VIOLET, size=14)
    line(s, x + 0.52, y + 0.02, 3, 0.3, "AI Case Assistant", 13, INK, bold=True)
    btn(s, x + w - 1.35, y + 0.03, "Generate", VIOLET)
    yy = y + 0.62
    # summary card
    rrect(s, x, yy, w, 1.35, VIOLET_LT, line=C('DDD6FE'), radius=0.08)
    line(s, x + 0.16, yy + 0.12, w, 0.3, "Case summary", 9.5, VIOLET, bold=True)
    for i, ww in enumerate([w - 0.4, w - 0.7, w - 1.3]):
        rrect(s, x + 0.16, yy + 0.46 + i * 0.24, ww, 0.12, C('DDD6FE'), radius=0.06)
    yy += 1.55
    line(s, x, yy, w, 0.3, "Suggested sections (IPC / BNS)", 9.5, INK, bold=True)
    cx = x
    for t, col, bgc in [("IPC 420", INDIGO, INDIGO_LT), ("BNS 318", INDIGO, INDIGO_LT),
                        ("IPC 406", INDIGO, INDIGO_LT), ("BNS 316", INDIGO, INDIGO_LT)]:
        cx += chip(s, cx, yy + 0.34, t, bgc, col, size=10, h=0.30) + 0.12
    yy += 0.95
    rrect(s, x, yy, w, 0.95, EMERALD_LT, line=C('A7F3D0'), radius=0.08)
    line(s, x + 0.16, yy + 0.13, 3, 0.3, "Suggested priority", 9.5, EMERALD_DK, bold=True)
    chip(s, x + 0.16, yy + 0.45, "HIGH", EMERALD, WHITE, size=10, h=0.30)
    line(s, x + w - 4.4, yy + 0.13, 4.2, 0.7,
         "Key facts extracted · disclaimer attached", 8.5, SLATE, align=PP_ALIGN.RIGHT)


def mk_crossexam(s, x, y, w, h):
    # tabs
    tw = (w - 0.12) / 2
    rrect(s, x, y, tw, 0.42, ROSE_LT, line=ROSE, radius=0.08)
    line(s, x, y, tw, 0.42, "⚔  Opponent", 11, ROSE_DK, bold=True, font=EMO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rrect(s, x + tw + 0.12, y, tw, 0.42, WHITE, line=LINE, radius=0.08)
    line(s, x + tw + 0.12, y, tw, 0.42, "⚖  Judge", 11, SLATE, bold=True, font=EMO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    yy = y + 0.6
    qs = [("Credibility", "Why was the FIR delayed by 11 days?"),
          ("Timeline", "Can you account for the gap in the evidence chain?"),
          ("Documentary", "Is there proof the payment was actually made?")]
    for cat, q in qs:
        rrect(s, x, yy, w, 0.92, WHITE, line=LINE, radius=0.08)
        chip(s, x + 0.16, yy + 0.13, cat, ROSE_LT, ROSE_DK, size=8.5, h=0.24)
        tb(s, x + 0.16, yy + 0.40, w - 0.3, 0.4,
           [P([R(q, 10.5, INK, True)])])
        line(s, x + 0.16, yy + 0.66, w - 0.3, 0.3,
             "Prep strategy: keep the timeline and documents ready.", 8.5, MUTED)
        yy += 1.04


def mk_memory(s, x, y, w, h):
    # amber staleness alert
    rrect(s, x, y, w, 0.92, AMBER_LT, line=AMBER, radius=0.08)
    icon_badge(s, x + 0.16, y + 0.22, 0.46, "⚠", AMBER, size=16)
    tb(s, x + 0.78, y + 0.15, w - 2.6, 0.7,
       [P([R("This case has changed since the AI notes were made.", 10.5, AMBER_DK, True)]),
        P([R("Regenerate to refresh the analysis.", 9.5, AMBER_DK)], space_before=2)])
    btn(s, x + w - 1.55, y + 0.28, "Regenerate", AMBER)
    yy = y + 1.2
    # cached state row
    rrect(s, x, yy, w, 0.78, EMERALD_LT, line=C('A7F3D0'), radius=0.08)
    icon_badge(s, x + 0.16, yy + 0.18, 0.42, "⚡", EMERALD, size=14)
    tb(s, x + 0.72, yy + 0.14, w - 1, 0.6,
       [P([R("Loaded instantly from saved results", 10.5, EMERALD_DK, True)]),
        P([R("No waiting, no re-charging when nothing changed", 9, EMERALD_DK)],
          space_before=2)])
    yy += 1.05
    # how it works mini diagram
    line(s, x, yy, w, 0.3, "How it stays fresh", 9.5, INK, bold=True)
    yy += 0.36
    steps = [("Case facts", SKY), ("+ timeline", INDIGO), ("= fingerprint", VIOLET)]
    sx = x
    sw = (w - 0.6) / 3
    for i, (t, col) in enumerate(steps):
        rrect(s, sx, yy, sw, 0.6, WHITE, line=col, radius=0.08)
        line(s, sx, yy, sw, 0.6, t, 9.5, col, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        if i < 2:
            line(s, sx + sw - 0.02, yy, 0.34, 0.6, "→", 16, MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        sx += sw + 0.3
    line(s, x, yy + 0.7, w, 0.4,
         "If the fingerprint changes, LexCase asks you to refresh.", 9, MUTED)


def mk_cases(s, x, y, w, h):
    line(s, x, y, 3, 0.3, "Cases", 12, INK, bold=True)
    btn(s, x + w - 1.45, y - 0.02, "+ New case", INDIGO)
    # search bar
    rrect(s, x, y + 0.42, w, 0.34, C('F1F5F9'), line=LINE, radius=0.10)
    line(s, x + 0.14, y + 0.40, 4, 0.34, "🔍  Search cases…", 9, MUTED, font=EMO,
         anchor=MSO_ANCHOR.MIDDLE)
    yy = y + 0.95
    rows = [("State vs. R. Kumar", "Criminal", "Trial", EMERALD, "Active"),
            ("Mehta Property Dispute", "Civil", "Hearing", AMBER, "Pending"),
            ("Sharma vs. Sharma", "Family", "Final args", EMERALD, "Active"),
            ("ACME Contract Breach", "Corporate", "Filed", SKY, "New")]
    # header
    line(s, x + 0.16, yy, 3, 0.25, "CASE", 8, MUTED, bold=True)
    line(s, x + 3.0, yy, 1.5, 0.25, "TYPE", 8, MUTED, bold=True)
    line(s, x + 4.5, yy, 1.5, 0.25, "STATUS", 8, MUTED, bold=True)
    yy += 0.32
    for nm, ty, stg, col, st in rows:
        rrect(s, x, yy, w, 0.62, WHITE, line=LINE, radius=0.07)
        tb(s, x + 0.16, yy + 0.10, 2.9, 0.5,
           [P([R(nm, 10, INK, True)]), P([R(stg, 8, MUTED)], space_before=1)])
        line(s, x + 3.0, yy, 1.5, 0.62, ty, 9, SLATE, anchor=MSO_ANCHOR.MIDDLE)
        chip(s, x + 4.5, yy + 0.18, st, INDIGO_LT, INDIGO, size=8.5, h=0.26)
        yy += 0.70


def mk_timeline(s, x, y, w, h):
    line(s, x, y, w, 0.3, "Case tracking timeline", 12, INK, bold=True)
    yy = y + 0.5
    stages = [("FIR / Complaint", "Sections noted", SKY),
              ("Investigation", "+ IPC 406 added", INDIGO),
              ("Charge Sheet", "Sections confirmed", VIOLET),
              ("Trial", "Evidence recorded", AMBER),
              ("Judgment", "Outcome logged", EMERALD)]
    lx = x + 0.35
    vline(s, lx, yy + 0.1, (len(stages) - 1) * 0.74 + 0.1, FAINT, 2.0)
    for i, (st, note, col) in enumerate(stages):
        cy = yy + i * 0.74
        oval(s, lx - 0.11, cy, 0.22, 0.22, col, line=WHITE, lw=2)
        tb(s, lx + 0.35, cy - 0.06, w - 1, 0.6,
           [P([R(st, 11, INK, True)]),
            P([R(note, 9, SLATE)], space_before=1)])
        if "+ " in note:
            chip(s, x + w - 1.7, cy, "section added", AMBER_LT, AMBER_DK, size=8, h=0.24)


def mk_favorability(s, x, y, w, h):
    line(s, x, y, w, 0.3, "Case favorability", 12, INK, bold=True,
         align=PP_ALIGN.CENTER)
    dd = min(w - 1.5, h - 2.0)
    dx = x + (w - dd) / 2
    dy = y + 0.55
    donut(s, dx, dy, dd, EMERALD, thickness=0.30)
    tb(s, dx, dy + dd / 2 - 0.45, dd, 1.0,
       [P([R("72", 40, INK, True)], PP_ALIGN.CENTER),
        P([R("out of 100", 11, MUTED)], PP_ALIGN.CENTER, space_before=2)],
       anchor=MSO_ANCHOR.MIDDLE)
    chip(s, x + w / 2 - 0.95, dy + dd + 0.18, "Strongly in your favour", EMERALD_LT,
         EMERALD_DK, size=10, h=0.32)
    line(s, x, dy + dd + 0.65, w, 0.5,
         "A simple score that shows, at a glance, how strong the matter looks.",
         9.5, MUTED, align=PP_ALIGN.CENTER)


def mk_clients(s, x, y, w, h):
    line(s, x, y, 3, 0.3, "Clients", 12, INK, bold=True)
    btn(s, x + w - 1.55, y - 0.02, "+ New client", SKY)
    yy = y + 0.55
    rows = [("Rahul Kumar", "Individual", "RK", INDIGO, "3 cases"),
            ("ACME Pvt. Ltd.", "Company", "AC", SKY, "2 cases"),
            ("Priya Mehta", "Individual", "PM", VIOLET, "1 case"),
            ("Sharma Family", "Individual", "SF", EMERALD, "1 case")]
    for nm, ty, ini, col, meta in rows:
        rrect(s, x, yy, w, 0.74, WHITE, line=LINE, radius=0.08)
        avatar(s, x + 0.16, yy + 0.16, 0.42, ini, col)
        tb(s, x + 0.72, yy + 0.14, 3.5, 0.5,
           [P([R(nm, 10.5, INK, True)]),
            P([R(ty, 8.5, MUTED)], space_before=1)])
        chip(s, x + w - 1.35, yy + 0.24, meta, C('F1F5F9'), SLATE, size=8.5, h=0.26)
        yy += 0.82


def mk_hearings(s, x, y, w, h):
    line(s, x, y, 3, 0.3, "Hearings calendar", 12, INK, bold=True)
    # mini calendar
    calw = w * 0.52
    yy = y + 0.5
    days = ["S", "M", "T", "W", "T", "F", "S"]
    cw = calw / 7
    for i, d in enumerate(days):
        line(s, x + i * cw, yy, cw, 0.24, d, 8, MUTED, bold=True, align=PP_ALIGN.CENTER)
    hl = {10, 17, 23}
    n = 1
    gy = yy + 0.3
    for r in range(4):
        for c in range(7):
            if n > 28:
                break
            cellx = x + c * cw
            celly = gy + r * 0.52
            if n in hl:
                rrect(s, cellx + 0.04, celly, cw - 0.08, 0.46, INDIGO, radius=0.08)
                line(s, cellx, celly, cw, 0.46, str(n), 9, WHITE, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            else:
                line(s, cellx, celly, cw, 0.46, str(n), 9, SLATE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            n += 1
    # agenda list
    ax = x + calw + 0.3
    aw = w - calw - 0.3
    line(s, ax, y + 0.5, aw, 0.24, "Upcoming", 9.5, INK, bold=True)
    ag = [("10", "State vs. Kumar", "Bail hearing", EMERALD),
          ("17", "Mehta Dispute", "Arguments", AMBER),
          ("23", "ACME Breach", "First hearing", SKY)]
    ay = y + 0.84
    for d, cs, pur, col in ag:
        rrect(s, ax, ay, aw, 0.78, WHITE, line=LINE, radius=0.08)
        rrect(s, ax + 0.12, ay + 0.14, 0.5, 0.5, col, radius=0.08)
        line(s, ax + 0.12, ay + 0.14, 0.5, 0.5, d, 13, WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tb(s, ax + 0.74, ay + 0.13, aw - 0.85, 0.55,
           [P([R(cs, 9.5, INK, True)]), P([R(pur, 8.5, MUTED)], space_before=1)])
        ay += 0.86


def mk_tasks(s, x, y, w, h):
    line(s, x, y, 3, 0.3, "Task board", 12, INK, bold=True)
    cols = [("To do", FAINT, [("Draft bail plea", ROSE), ("Call client", AMBER)]),
            ("In progress", AMBER, [("Collect evidence", AMBER), ("File reply", INDIGO)]),
            ("Done", EMERALD, [("Notice served", EMERALD)])]
    yy = y + 0.5
    cw = (w - 0.4) / 3
    for i, (title, col, cards) in enumerate(cols):
        cx = x + i * (cw + 0.2)
        rrect(s, cx, yy, cw, h - 0.6, C('F1F5F9'), radius=0.08)
        oval(s, cx + 0.14, yy + 0.16, 0.12, 0.12, col)
        line(s, cx + 0.34, yy + 0.08, cw, 0.26, title, 9.5, INK, bold=True)
        ty = yy + 0.5
        for name, pcol in cards:
            rrect(s, cx + 0.12, ty, cw - 0.24, 0.72, WHITE, line=LINE, radius=0.07)
            rect(s, cx + 0.12, ty + 0.07, 0.05, 0.58, pcol)
            tb(s, cx + 0.28, ty + 0.10, cw - 0.42, 0.55,
               [P([R(name, 9, INK, True)], line=1.05),
                P([R("Due in 2 days", 7.5, MUTED)], space_before=2)])
            ty += 0.82


def mk_docs(s, x, y, w, h):
    line(s, x, y, 4, 0.3, "Documents", 12, INK, bold=True)
    yy = y + 0.45
    docs = [("Vakalatnama.pdf", "v3", INDIGO), ("Bail application.docx", "v2", SKY),
            ("Evidence list.pdf", "v1", VIOLET)]
    for nm, ver, col in docs:
        rrect(s, x, yy, w, 0.56, WHITE, line=LINE, radius=0.07)
        icon_badge(s, x + 0.12, yy + 0.11, 0.34, "📄", col, size=11)
        line(s, x + 0.6, yy, 4, 0.56, nm, 9.5, INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        chip(s, x + w - 0.85, yy + 0.16, ver, INDIGO_LT, INDIGO, size=8.5, h=0.24)
        yy += 0.60
    yy += 0.10
    line(s, x, yy, 4, 0.3, "Evidence locker", 12, INK, bold=True)
    yy += 0.42
    ev = [("Bank statement", "Documentary", "Verified", EMERALD),
          ("CCTV footage", "Electronic", "In custody", AMBER)]
    for nm, ty, st, col in ev:
        rrect(s, x, yy, w, 0.62, WHITE, line=LINE, radius=0.07)
        icon_badge(s, x + 0.12, yy + 0.13, 0.36, "🔏", col, size=11)
        tb(s, x + 0.62, yy + 0.12, 4, 0.5,
           [P([R(nm, 9.5, INK, True)]), P([R(ty, 8, MUTED)], space_before=1)])
        chip(s, x + w - 1.25, yy + 0.18, st, EMERALD_LT if col == EMERALD else AMBER_LT,
             EMERALD_DK if col == EMERALD else AMBER_DK, size=8, h=0.26)
        yy += 0.70


def mk_library(s, x, y, w, h):
    line(s, x, y, 5, 0.3, "Legal Library — ready-made templates", 11.5, INK, bold=True)
    yy = y + 0.46
    tmpl = [("Bail Application", "Criminal"), ("Vakalatnama", "General"),
            ("Legal Notice", "Civil"), ("Rent Agreement", "Property")]
    cw = (w - 0.2) / 2
    for i, (nm, cat) in enumerate(tmpl):
        cx = x + (i % 2) * (cw + 0.2)
        cyy = yy + (i // 2) * 0.82
        rrect(s, cx, cyy, cw, 0.72, WHITE, line=LINE, radius=0.08)
        icon_badge(s, cx + 0.12, cyy + 0.16, 0.4, "📜", VIOLET, size=12)
        tb(s, cx + 0.62, cyy + 0.13, cw - 0.7, 0.5,
           [P([R(nm, 9.5, INK, True)]), P([R(cat + " · editable", 8, MUTED)],
                                          space_before=1)])
    yy += 1.85
    line(s, x, yy, 5, 0.3, "Legal Notebook — quick law reference", 11.5, INK, bold=True)
    yy += 0.42
    secs = [("IPC 420", "Cheating"), ("BNS 318", "Cheating (new code)"),
            ("IPC 302", "Murder")]
    for code, desc in secs:
        rrect(s, x, yy, w, 0.42, C('F8FAFC'), line=LINE, radius=0.06)
        chip(s, x + 0.12, yy + 0.08, code, INDIGO_LT, INDIGO, size=8.5, h=0.26)
        line(s, x + 1.4, yy, 4, 0.42, desc, 9, SLATE, anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.5


def mk_search(s, x, y, w, h):
    rrect(s, x + 0.3, y + 0.12, w - 0.6, h - 0.24, C('F1F5F9'), radius=0.10)
    px, py, pw = x + 0.6, y + 0.35, w - 1.2
    rrect(s, px, py, pw, h - 0.6, WHITE, line=INDIGO, lw=1.5, radius=0.10)
    rrect(s, px + 0.18, py + 0.16, pw - 0.36, 0.44, C('F8FAFC'), line=LINE, radius=0.10)
    line(s, px + 0.34, py + 0.14, pw, 0.44, "🔍   kumar", 13, INK, font=EMO,
         anchor=MSO_ANCHOR.MIDDLE)
    chip(s, px + pw - 0.82, py + 0.25, "⌘ K", INDIGO_LT, INDIGO, size=9, h=0.26, font=EMO)
    yy = py + 0.78
    groups = [("CASES", [("State vs. R. Kumar", "📂"), ("Kumar bail matter", "📂")]),
              ("CLIENTS", [("Rahul Kumar", "👥")]),
              ("HEARINGS", [("Kumar — bail hearing, Jun 10", "📅")])]
    for gname, items in groups:
        line(s, px + 0.22, yy, 4, 0.22, gname, 8, MUTED, bold=True)
        yy += 0.26
        for txt, g in items:
            rrect(s, px + 0.18, yy, pw - 0.36, 0.38, C('F8FAFC'), radius=0.07)
            line(s, px + 0.34, yy, 0.4, 0.38, g, 11, INK, font=EMO,
                 anchor=MSO_ANCHOR.MIDDLE)
            line(s, px + 0.74, yy, pw - 1, 0.38, txt, 9.5, INK, anchor=MSO_ANCHOR.MIDDLE)
            yy += 0.42
        yy += 0.05


def mk_team(s, x, y, w, h):
    line(s, x, y, 3, 0.3, "Team & roles", 12, INK, bold=True)
    yy = y + 0.5
    mem = [("Anita Rao", "Firm Owner", "AR", INDIGO),
           ("Priya Nair", "Partner", "PN", VIOLET),
           ("Rohan Das", "Associate", "RD", SKY),
           ("Sara Khan", "Paralegal", "SK", EMERALD)]
    for nm, role, ini, col in mem:
        rrect(s, x, yy, w * 0.62, 0.66, WHITE, line=LINE, radius=0.08)
        avatar(s, x + 0.14, yy + 0.13, 0.4, ini, col)
        tb(s, x + 0.66, yy + 0.12, 3, 0.5,
           [P([R(nm, 10, INK, True)]), P([R(role, 8.5, MUTED)], space_before=1)])
        yy += 0.74
    # permission matrix
    mx = x + w * 0.62 + 0.25
    mw = w - w * 0.62 - 0.25
    line(s, mx, y + 0.5, mw, 0.24, "Who can do what", 9.5, INK, bold=True)
    perms = [("View cases", True, True, True), ("Edit cases", True, True, False),
             ("Manage team", True, False, False), ("Delete", True, False, False)]
    ry = y + 0.84
    colx = [mx + 1.35, mx + 1.95, mx + 2.55]
    for lab, x1 in zip(["Own", "Par", "Ass"], colx):
        line(s, x1 - 0.2, ry - 0.26, 0.5, 0.24, lab, 7.5, MUTED, bold=True,
             align=PP_ALIGN.CENTER)
    for lab, a, b, c in perms:
        line(s, mx, ry, 1.4, 0.3, lab, 8.5, SLATE)
        for val, x1 in zip([a, b, c], colx):
            g = "✓" if val else "—"
            col = EMERALD if val else FAINT
            line(s, x1 - 0.2, ry - 0.02, 0.5, 0.3, g, 11, col, bold=True,
                 align=PP_ALIGN.CENTER)
        ry += 0.42


def mk_activity(s, x, y, w, h):
    line(s, x, y, 4, 0.3, "Activity log", 12, INK, bold=True)
    yy = y + 0.46
    acts = [("Priya", "created case", "ACME Breach", "2m", VIOLET),
            ("Rohan", "added hearing", "Jun 17", "1h", SKY),
            ("Anita", "updated client", "Rahul Kumar", "3h", INDIGO),
            ("Sara", "uploaded document", "Bail plea v2", "1d", EMERALD)]
    lx = x + 0.2
    vline(s, lx, yy + 0.15, (len(acts) - 1) * 0.62 + 0.1, FAINT, 1.5)
    for who, did, what, when, col in acts:
        oval(s, lx - 0.07, yy + 0.06, 0.16, 0.16, col, line=WHITE, lw=2)
        tb(s, lx + 0.3, yy, w - 1, 0.5,
           [P([R(who + " ", 9.5, INK, True), R(did + " ", 9.5, SLATE),
               R(what, 9.5, INK, True)])])
        line(s, x + w - 0.7, yy, 0.6, 0.3, when, 8, MUTED, align=PP_ALIGN.RIGHT)
        yy += 0.62
    yy += 0.1
    rrect(s, x, yy, w, 0.7, INDIGO_LT, radius=0.08)
    icon_badge(s, x + 0.14, yy + 0.16, 0.4, "🔔", INDIGO, size=13)
    tb(s, x + 0.66, yy + 0.13, w - 1, 0.5,
       [P([R("Notifications", 10, INK, True)]),
        P([R("You were assigned to “ACME Breach”", 8.5, SLATE)], space_before=1)])
    chip(s, x + w - 0.85, yy + 0.22, "New", INDIGO, WHITE, size=8, h=0.26)


def mk_security(s, x, y, w, h):
    line(s, x, y, w, 0.3, "Every firm's data stays private", 12, INK, bold=True,
         align=PP_ALIGN.CENTER)
    bw = (w - 1.4) / 2
    yy = y + 0.6
    bh = h - 1.5
    rrect(s, x, yy, bw, bh, INDIGO_LT, line=INDIGO, radius=0.10)
    line(s, x, yy + 0.18, bw, 0.3, "Firm A", 11, INDIGO_DK, bold=True,
         align=PP_ALIGN.CENTER)
    for i in range(3):
        rrect(s, x + 0.25, yy + 0.6 + i * 0.5, bw - 0.5, 0.36, WHITE, line=C('C7D2FE'),
              radius=0.06)
        line(s, x + 0.4, yy + 0.6 + i * 0.5, bw, 0.36,
             ["Clients", "Cases", "Documents"][i], 9, INDIGO, anchor=MSO_ANCHOR.MIDDLE)
    bx = x + bw + 1.4
    rrect(s, bx, yy, bw, bh, EMERALD_LT, line=EMERALD, radius=0.10)
    line(s, bx, yy + 0.18, bw, 0.3, "Firm B", 11, EMERALD_DK, bold=True,
         align=PP_ALIGN.CENTER)
    for i in range(3):
        rrect(s, bx + 0.25, yy + 0.6 + i * 0.5, bw - 0.5, 0.36, WHITE, line=C('A7F3D0'),
              radius=0.06)
        line(s, bx + 0.4, yy + 0.6 + i * 0.5, bw, 0.36,
             ["Clients", "Cases", "Documents"][i], 9, EMERALD_DK,
             anchor=MSO_ANCHOR.MIDDLE)
    # shield divider
    midx = x + bw + 0.45
    oval(s, midx, yy + bh / 2 - 0.35, 0.7, 0.7, WHITE, line=SLATE, lw=1.5)
    line(s, midx, yy + bh / 2 - 0.36, 0.7, 0.7, "🛡", 22, INDIGO, font=EMO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    line(s, x, yy + bh + 0.18, w, 0.4,
         "A firm can never see another firm's records — guaranteed by design.",
         9.5, MUTED, align=PP_ALIGN.CENTER)


# ============================================================================= TITLE SLIDE
def title_slide():
    s = slide()
    bg(s, C('0B1020'))
    # decorative glows
    oval(s, -2.0, -2.2, 6, 6, C('1E1B4B'))
    oval(s, 9.5, 3.8, 6.5, 6.5, C('15203B'))
    rect(s, 0, 0, 0.18, SH, INDIGO)
    # brand
    icon_badge(s, 0.95, 0.85, 0.7, "⚖", INDIGO, size=26)
    line(s, 1.8, 0.92, 6, 0.5, "LexCase", 26, WHITE, bold=True)
    line(s, 1.8, 1.42, 8, 0.3, "AI-Assisted Legal Case Management", 12, C('A5B4FC'))
    # headline
    tb(s, 0.95, 2.7, 9.6, 2.0,
       [P([R("Run your entire law practice", 40, WHITE, True)], line=1.05),
        P([R("from ", 40, WHITE, True), R("one smart workspace", 40, C('A5B4FC'), True)],
          line=1.05, space_before=4)])
    tb(s, 0.98, 4.55, 9.0, 0.8,
       [P([R("Cases, clients, hearings, tasks and documents — with an AI assistant "
            "that reads the case and helps you prepare.", 14, C('CBD5E1'))], line=1.3)])
    # feature pills
    px = 0.98
    for t, g in [("AI Case Assistant", "✨"), ("Cross-exam prep", "⚖"),
                 ("Smart dashboard", "📊"), ("All-in-one", "🗂")]:
        w = 0.4 + len(t) * 0.085 + 0.45
        rrect(s, px, 5.55, w, 0.5, C('1E2540'), line=C('334155'), radius=0.25)
        line(s, px + 0.22, 5.53, 0.4, 0.5, g, 13, WHITE, font=EMO,
             anchor=MSO_ANCHOR.MIDDLE)
        line(s, px + 0.62, 5.53, w, 0.5, t, 11, C('E2E8F0'), bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
        px += w + 0.25
    line(s, 0.98, 6.6, 11, 0.3,
         "Built for Indian lawyers & law firms  ·  A guided product tour", 11, C('64748B'))
    line(s, SW - 1.4, SH - 0.5, 0.7, 0.3, "01", 9, C('475569'), align=PP_ALIGN.RIGHT)
    _page['n'] = 1


# ============================================================================= OVERVIEW SLIDE
def overview_slide():
    s = slide()
    bg(s)
    rect(s, 0, 0, SW, 0.16, INDIGO)
    line(s, 0.7, 0.55, 8, 0.3, "WHAT IS LEXCASE?", 12, INDIGO, bold=True)
    tb(s, 0.7, 0.95, 11.8, 1.0,
       [P([R("One place to manage every case — start to finish", 28, INK, True)])])
    tb(s, 0.7, 1.85, 11.6, 0.7,
       [P([R("LexCase keeps a law firm's clients, cases, hearings, tasks, documents and "
            "team together in a single, easy workspace — and adds an AI assistant that "
            "helps lawyers prepare faster.", 14, SLATE)], line=1.3)])
    # 3 value cards
    cards = [("🗂", "Everything together", "No more scattered files, sheets and "
              "WhatsApp notes. The whole matter lives in one record.", INDIGO, INDIGO_LT),
             ("🤖", "AI that helps", "It reads the case and drafts summaries, suggests "
              "the right legal sections and prepares cross-examination questions.", VIOLET,
              VIOLET_LT),
             ("🔐", "Safe & private", "Each firm's data is fully separated, with clear "
              "roles deciding who can see and do what.", EMERALD, EMERALD_LT)]
    cw = (SW - 1.4 - 0.6) / 3
    for i, (g, t, d, col, bgc) in enumerate(cards):
        cx = 0.7 + i * (cw + 0.3)
        rrect(s, cx, 2.95, cw, 2.55, WHITE, line=LINE, radius=0.10)
        rect(s, cx, 2.95, cw, 0.10, col)
        icon_badge(s, cx + 0.3, 3.3, 0.62, g, col, size=22)
        tb(s, cx + 0.3, 4.1, cw - 0.6, 0.5, [P([R(t, 15, INK, True)])])
        tb(s, cx + 0.3, 4.55, cw - 0.6, 1.0, [P([R(d, 11, SLATE)], line=1.25)])
    # stat strip
    stats = [("12", "modules"), ("1", "AI assistant"), ("5", "firm roles"),
             ("100%", "data isolation")]
    sx = 0.7
    sw = (SW - 1.4 - 0.9) / 4
    for num, lab in stats:
        rrect(s, sx, 5.75, sw, 1.0, INDIGO_LT, radius=0.10)
        line(s, sx, 5.9, sw, 0.5, num, 24, INDIGO_DK, bold=True, align=PP_ALIGN.CENTER)
        line(s, sx, 6.42, sw, 0.3, lab, 10, SLATE, align=PP_ALIGN.CENTER)
        sx += sw + 0.3
    footer(s)


# ============================================================================= PROBLEM SLIDE
def problem_slide():
    s = slide()
    bg(s)
    rect(s, 0, 0, SW, 0.16, ROSE)
    line(s, 0.7, 0.55, 8, 0.3, "THE CHALLENGE", 12, ROSE_DK, bold=True)
    tb(s, 0.7, 0.95, 11.8, 1.0,
       [P([R("Running a law practice is a juggling act", 28, INK, True)])])
    # left: scattered problems
    probs = [("📑", "Files everywhere", "Case papers spread across folders, email and chat."),
             ("⏰", "Missed dates", "Hearing dates and deadlines easy to forget."),
             ("🔁", "Repeated work", "The same summaries and notices typed again and again."),
             ("❓", "Hard to see status", "No quick way to know how a case is really doing.")]
    yy = 2.1
    for g, t, d in probs:
        rrect(s, 0.7, yy, 5.6, 0.95, WHITE, line=LINE, radius=0.10)
        icon_badge(s, 0.92, yy + 0.22, 0.5, g, ROSE_LT, ROSE_DK, size=18)
        tb(s, 1.6, yy + 0.18, 4.5, 0.7,
           [P([R(t, 13, INK, True)]),
            P([R(d, 10.5, SLATE)], space_before=2)])
        yy += 1.08
    # right: the answer
    rrect(s, 6.7, 2.1, 5.9, 4.2, INDIGO, radius=0.12)
    icon_badge(s, 7.05, 2.45, 0.7, "💡", WHITE, INDIGO, size=24)
    tb(s, 7.05, 3.35, 5.2, 0.6, [P([R("LexCase brings it all together", 19, WHITE, True)],
                                    line=1.05)])
    pts = ["One record per case — papers, dates, tasks, notes",
           "Reminders and a calendar so nothing slips",
           "Ready-made templates instead of retyping",
           "A clear score and dashboard for instant status",
           "An AI assistant that does the heavy reading"]
    py = 4.05
    for p in pts:
        line(s, 7.05, py, 0.4, 0.3, "✓", 13, C('A5B4FC'), bold=True)
        tb(s, 7.4, py, 4.9, 0.5, [P([R(p, 11.5, WHITE)], line=1.1)])
        py += 0.45
    footer(s)


# ============================================================================= APP MAP SLIDE
def appmap_slide():
    s = slide()
    bg(s)
    rect(s, 0, 0, SW, 0.16, INDIGO)
    line(s, 0.7, 0.55, 8, 0.3, "EVERYTHING IN ONE PLACE", 12, INDIGO, bold=True)
    tb(s, 0.7, 0.95, 11.8, 0.8, [P([R("The complete workspace, at a glance", 28, INK, True)])])
    line(s, 0.7, 1.75, 11.6, 0.4,
         "Simple menus group every tool a firm needs. Here's the full map.", 13, SLATE)
    groups = [
        ("WORKSPACE", INDIGO, [("📊", "Dashboard")]),
        ("MATTERS", VIOLET, [("📂", "Cases"), ("👥", "Clients"),
                             ("📅", "Hearings"), ("✅", "Tasks")]),
        ("RECORDS", SKY, [("📁", "Documents"), ("🔏", "Evidence"),
                          ("📜", "Legal Library"), ("📖", "Legal Notebook")]),
        ("FIRM", EMERALD, [("👤", "Team"), ("🔑", "Roles & Rights"),
                           ("🛡", "Activity Log")]),
    ]
    colw = (SW - 1.4 - 0.9) / 4
    cx = 0.7
    for gname, col, items in groups:
        rrect(s, cx, 2.35, colw, 4.3, WHITE, line=LINE, radius=0.10)
        rect(s, cx, 2.35, colw, 0.5, col)
        line(s, cx, 2.4, colw, 0.4, gname, 11, WHITE, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        iy = 3.05
        for g, name in items:
            rrect(s, cx + 0.2, iy, colw - 0.4, 0.72, BG, line=LINE, radius=0.08)
            icon_badge(s, cx + 0.34, iy + 0.16, 0.4, g, col, size=13)
            line(s, cx + 0.85, iy, colw - 1.05, 0.72, name, 11, INK, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
            iy += 0.85
        cx += colw + 0.3
    footer(s)


# ============================================================================= SECTION DIVIDER
def divider(text, sub, glyph, color):
    s = slide()
    bg(s, C('0B1020'))
    oval(s, -2, 3, 6, 6, C('15203B'))
    oval(s, 9.5, -2.5, 6.5, 6.5, C('1E1B4B'))
    rect(s, 0, 0, 0.18, SH, color)
    icon_badge(s, 0.95, 2.7, 0.9, glyph, color, size=34)
    tb(s, 0.95, 3.85, 11, 1.0, [P([R(text, 34, WHITE, True)])])
    tb(s, 0.98, 4.8, 10.5, 0.8, [P([R(sub, 14, C('CBD5E1'))], line=1.3)])
    line(s, SW - 1.4, SH - 0.5, 0.7, 0.3, str(_page['n'] + 1).zfill(2), 9, C('475569'),
         align=PP_ALIGN.RIGHT)
    _page['n'] += 1


# ============================================================================= CLOSING / WHY
def why_slide():
    s = slide()
    bg(s)
    rect(s, 0, 0, SW, 0.16, INDIGO)
    line(s, 0.7, 0.55, 8, 0.3, "WHY LEXCASE", 12, INDIGO, bold=True)
    tb(s, 0.7, 0.95, 11.8, 0.8, [P([R("Less admin, more lawyering", 28, INK, True)])])
    benefits = [
        ("⏱", "Save hours every week", "Templates, search and AI summaries cut the "
         "repetitive work dramatically.", INDIGO, INDIGO_LT),
        ("🎯", "Be better prepared", "AI cross-exam questions and section suggestions "
         "help you walk in ready.", VIOLET, VIOLET_LT),
        ("📈", "See the whole picture", "Dashboard, win-rate and favorability scores "
         "show status at a glance.", SKY, SKY_LT),
        ("🤝", "Work as a team", "Shared cases, clear roles and an activity trail keep "
         "everyone aligned.", EMERALD, EMERALD_LT),
        ("🔒", "Trust your data", "Each firm's information is isolated and access is "
         "tightly controlled.", AMBER, AMBER_LT),
        ("📱", "Works anywhere", "A clean, modern interface that works on desktop and "
         "mobile.", ROSE, ROSE_LT),
    ]
    cw = (SW - 1.4 - 0.6) / 3
    ch = 1.95
    for i, (g, t, d, col, bgc) in enumerate(benefits):
        cx = 0.7 + (i % 3) * (cw + 0.3)
        cy = 2.0 + (i // 3) * (ch + 0.3)
        rrect(s, cx, cy, cw, ch, WHITE, line=LINE, radius=0.10)
        icon_badge(s, cx + 0.28, cy + 0.28, 0.56, g, col, size=20)
        tb(s, cx + 0.28, cy + 1.0, cw - 0.56, 0.4, [P([R(t, 14, INK, True)])])
        tb(s, cx + 0.28, cy + 1.4, cw - 0.56, 0.6, [P([R(d, 10.5, SLATE)], line=1.2)])
    footer(s)


def thanks_slide():
    s = slide()
    bg(s, C('0B1020'))
    oval(s, -2.5, -2.5, 7, 7, C('1E1B4B'))
    oval(s, 9, 4, 6.5, 6.5, C('15203B'))
    rect(s, 0, 0, 0.18, SH, INDIGO)
    icon_badge(s, 0.95, 1.6, 0.7, "⚖", INDIGO, size=26)
    line(s, 1.8, 1.7, 6, 0.5, "LexCase", 22, WHITE, bold=True)
    tb(s, 0.95, 2.9, 11, 1.5, [P([R("Thank you", 46, WHITE, True)])])
    tb(s, 0.98, 4.2, 10, 0.8,
       [P([R("One smart workspace for the whole practice — powered by AI.", 15,
            C('CBD5E1'))], line=1.3)])
    # CTA pills
    px = 0.98
    for t, g in [("Request a demo", "▶"), ("Try the sample firm", "🧪")]:
        w = 0.5 + len(t) * 0.1 + 0.5
        rrect(s, px, 5.3, w, 0.6, INDIGO if px < 2 else C('1E2540'),
              line=None if px < 2 else C('334155'), radius=0.3)
        line(s, px + 0.28, 5.28, 0.4, 0.6, g, 13, WHITE, font=EMO,
             anchor=MSO_ANCHOR.MIDDLE)
        line(s, px + 0.68, 5.28, w, 0.6, t, 12, WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
        px += w + 0.3
    line(s, 0.98, 6.5, 11, 0.3,
         "Disclaimer: AI features assist a qualified lawyer — they are not legal advice.",
         10, C('64748B'))
    line(s, SW - 1.4, SH - 0.5, 0.7, 0.3, str(_page['n'] + 1).zfill(2), 9, C('475569'),
         align=PP_ALIGN.RIGHT)


# ============================================================================= BUILD
title_slide()
overview_slide()
problem_slide()
appmap_slide()

divider("Your daily command center", "See the health of the whole practice the moment "
        "you log in.", "📊", SKY)

feature(1, "Dashboard", "📊", SKY, SKY_LT,
        "A live dashboard of your whole firm",
        "The moment you log in, see the numbers that matter — without opening a single "
        "case.",
        [("Key numbers —", "active cases, clients, open tasks and this week's hearings."),
         ("Trends —", "a 6-month view of cases and hearings at a glance."),
         ("Win rate —", "how many matters the firm is winning."),
         ("Workload —", "see how work is shared across each lawyer.")],
        mk_dashboard, "lexcase.app/dashboard")

divider("Meet your AI legal assistant", "It reads the case and does the heavy "
        "preparation with you.", "🤖", VIOLET)

feature(2, "AI Case Assistant", "✨", VIOLET, VIOLET_LT,
        "AI that reads the case for you",
        "Feed it the case facts and it writes a clean summary and points you to the "
        "likely legal sections.",
        [("Plain summary —", "3–5 sentences plus the key facts, instantly."),
         ("Right sections —", "suggests the IPC / BNS sections that likely apply."),
         ("Smart priority —", "recommends how urgent the matter is."),
         ("Always safe —", "every suggestion is flagged for a lawyer to review.")],
        mk_ai, "lexcase.app/cases/ai")

feature(3, "Cross-Examination Prep", "⚖", ROSE, ROSE_LT,
        "Walk in ready for the hearing",
        "LexCase predicts the tough questions — from both the opposing counsel and the "
        "judge — before you reach court.",
        [("Opponent's angle —", "questions made to attack credibility and find gaps."),
         ("Judge's angle —", "questions to test the facts and the legal basis."),
         ("Categorised —", "each question is tagged (credibility, timeline, …)."),
         ("Prep strategy —", "a short tip on how to answer each one.")],
        mk_crossexam, "lexcase.app/cases/cross-exam")

feature(4, "Smart AI Memory", "⚡", AMBER, AMBER_LT,
        "Instant answers that stay fresh",
        "AI results are saved with each case, so revisiting is instant — and LexCase "
        "tells you when they need a refresh.",
        [("Instant on return —", "saved results open immediately, no waiting."),
         ("No wasted effort —", "it won't re-run when nothing has changed."),
         ("Knows when stale —", "spots when the case has changed since last time."),
         ("One-click refresh —", "a clear ‘Regenerate' button updates everything.")],
        mk_memory, "lexcase.app/cases")

divider("Manage every matter end-to-end", "From the first complaint to the final "
        "judgment — all in one record.", "📂", INDIGO)

feature(5, "Case Management", "📂", INDIGO, INDIGO_LT,
        "Every case, fully organised",
        "Each case ties together the client, lawyers, hearings, tasks, documents and "
        "notes — in one tidy place.",
        [("One record —", "client, court, team and papers all in one view."),
         ("Search & filter —", "find any matter in seconds."),
         ("Status at a glance —", "clear badges for stage, type and priority."),
         ("Safe archive —", "archive and restore instead of losing anything.")],
        mk_cases, "lexcase.app/cases")

feature(6, "Case Tracking Timeline", "🧭", VIOLET, VIOLET_LT,
        "Follow a case, stage by stage",
        "Track the matter as it moves — FIR to investigation, charges, trial and "
        "judgment — with what changed at each step.",
        [("Full journey —", "every stage of the case in order."),
         ("Sections snapshot —", "see which legal sections applied at each stage."),
         ("Highlights changes —", "new sections added are clearly marked."),
         ("AI help —", "it can even suggest sections for a new update.")],
        mk_timeline, "lexcase.app/cases/timeline")

feature(7, "Favorability Score", "🎯", EMERALD, EMERALD_LT,
        "How strong is the case? One number.",
        "A simple 0–100 gauge shows, at a glance, how strongly the matter is leaning in "
        "your favour.",
        [("Instant read —", "no spreadsheets, just one clear score."),
         ("Visual gauge —", "a colour-coded dial anyone can understand."),
         ("Smarter decisions —", "spot strong and weak matters quickly."),
         ("Client-friendly —", "easy to explain progress to clients.")],
        mk_favorability, "lexcase.app/cases/show")

divider("Clients, dates and to-dos", "Keep people, hearings and daily work moving "
        "without dropping a thing.", "👥", SKY)

feature(8, "Clients", "👥", SKY, SKY_LT,
        "Every client, neatly on file",
        "Keep full client profiles — individuals or companies — linked straight to their "
        "cases.",
        [("Complete profiles —", "contact details and type in one card."),
         ("Linked to cases —", "jump from a client to all their matters."),
         ("Quick add & edit —", "simple forms, no clutter."),
         ("Easy to find —", "search clients instantly.")],
        mk_clients, "lexcase.app/clients")

feature(9, "Hearings", "📅", AMBER, AMBER_LT,
        "Never miss a court date again",
        "A clean calendar and agenda of every hearing, with purpose, judge and outcome "
        "tracked.",
        [("Calendar view —", "see all upcoming dates at a glance."),
         ("Full details —", "purpose, judge and status for each hearing."),
         ("Outcome tracking —", "record what happened after each date."),
         ("Stay ahead —", "an agenda of what's coming up next.")],
        mk_hearings, "lexcase.app/hearings")

feature(10, "Tasks", "✅", EMERALD, EMERALD_LT,
        "A clear to-do board for the team",
        "Organise work on a simple board — drag tasks across columns, set due dates and "
        "track who's doing what.",
        [("Visual board —", "To-do, In-progress and Done at a glance."),
         ("Drag to update —", "move a task as work progresses."),
         ("Due dates & priority —", "know what's urgent."),
         ("Full history —", "see how each task moved over time.")],
        mk_tasks, "lexcase.app/tasks")

divider("All your paperwork, sorted", "Documents, evidence and ready-made legal "
        "templates — organised and reusable.", "🗂", VIOLET)

feature(11, "Documents & Evidence", "📁", VIOLET, VIOLET_LT,
        "Papers and proof, properly tracked",
        "Keep documents organised with versions, and log evidence with its type and "
        "chain-of-custody status.",
        [("Versioned files —", "always know which draft is the latest."),
         ("Tidy folders —", "documents organised per case."),
         ("Evidence locker —", "log each item with its type."),
         ("Chain of custody —", "track an item's status from intake to court.")],
        mk_docs, "lexcase.app/documents")

feature(12, "Legal Library & Notebook", "📜", INDIGO, INDIGO_LT,
        "Templates and law, ready to use",
        "A library of editable document templates, plus a quick-reference notebook of "
        "Indian statutes and sections.",
        [("Ready templates —", "bail pleas, notices, agreements and more."),
         ("Edit & reuse —", "duplicate and customise in seconds."),
         ("Printable —", "produce clean documents fast."),
         ("Law at hand —", "look up sections without leaving the app.")],
        mk_library, "lexcase.app/templates")

feature(13, "Quick Search", "🔍", SKY, SKY_LT,
        "Find anything in one keystroke",
        "Press ⌘K and search the whole firm — cases, clients, hearings and tasks — as "
        "you type.",
        [("One search box —", "covers every part of the app."),
         ("As-you-type —", "results appear instantly."),
         ("Neatly grouped —", "by cases, clients, hearings and tasks."),
         ("Jump anywhere —", "open the right record in a click.")],
        mk_search, "lexcase.app  ·  press ⌘K")

divider("Run the firm, safely", "The right people, the right access, and a full record "
        "of who did what.", "🔐", EMERALD)

feature(14, "Team & Access Control", "👤", EMERALD, EMERALD_LT,
        "The right access for every person",
        "Add your team and give each person a role — from Firm Owner to Clerk — that "
        "decides exactly what they can see and do.",
        [("Five clear roles —", "Owner, Partner, Associate, Paralegal, Clerk."),
         ("Fine-grained rights —", "control each ability per role."),
         ("Add & manage —", "invite, edit or remove members easily."),
         ("Right-sized access —", "people only see what they should.")],
        mk_team, "lexcase.app/team")

feature(15, "Activity Log & Alerts", "🛡", INDIGO, INDIGO_LT,
        "A clear record of every action",
        "See who changed what and when across the firm, and keep everyone informed with "
        "in-app notifications.",
        [("Full audit trail —", "every important change is logged."),
         ("Accountability —", "know exactly who did what, when."),
         ("In-app alerts —", "get notified when assigned to a matter."),
         ("Stay in the loop —", "a tidy notification centre, mark-all-read.")],
        mk_activity, "lexcase.app/activity")

feature(16, "Privacy & Security", "🔒", AMBER, AMBER_LT,
        "Your firm's data stays your firm's",
        "LexCase is multi-firm by design — every firm's records are completely separated, "
        "so data can never leak across firms.",
        [("Total separation —", "one firm can't see another's data."),
         ("Protected actions —", "every action is permission-checked."),
         ("Hidden IDs —", "records use safe, non-guessable links."),
         ("Audited —", "a full trail backs up every change.")],
        mk_security, "lexcase.app")

why_slide()
thanks_slide()

# ----------------------------------------------------------------------------- save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "LexCase-Features-Overview.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
